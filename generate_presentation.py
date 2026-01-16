"""
Generate PowerPoint Presentation for NIFTY Quantitative Trading Strategy
Updated with latest LSTM improvements and results
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Professional dark theme
DARK_BLUE = RGBColor(15, 32, 65)
ACCENT_BLUE = RGBColor(0, 122, 204)
GREEN = RGBColor(40, 180, 99)
RED = RGBColor(231, 76, 60)
GRAY = RGBColor(128, 128, 128)
LIGHT_GRAY = RGBColor(200, 200, 200)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = ACCENT_BLUE
        p2.alignment = PP_ALIGN.CENTER
    
    # Author line
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Kartik Mendiratta | January 2026"
    p3.font.size = Pt(16)
    p3.font.color.rgb = LIGHT_GRAY
    p3.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Bullet points
    has_image = image_path and os.path.exists(image_path)
    if has_image:
        left = Inches(0.5)
        width = Inches(5.8)
        font_size = Pt(16)
        space_after = Pt(10)
    else:
        left = Inches(0.5)
        width = Inches(12)
        font_size = Pt(18)
        space_after = Pt(14)
    
    txBox2 = slide.shapes.add_textbox(left, Inches(1.5), width, Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "▸ " + bullet
        p.font.size = font_size
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = space_after
    
    # Image
    if has_image:
        slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), width=Inches(6.3))
    
    # Slide number
    add_slide_number(slide, len(prs.slides))
    return slide

def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_section_header(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2.8), Inches(3.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    add_slide_number(slide, len(prs.slides))
    return slide

def add_results_table_slide(title, headers, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Create table
    rows = len(data) + 1
    cols = len(headers)
    table_width = Inches(10)
    table_height = Inches(4)
    
    x = (prs.slide_width - table_width) / 2
    table = slide.shapes.add_table(rows, cols, x, Inches(2), table_width, table_height).table
    
    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        tf = cell.text_frame
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            tf = cell.text_frame
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Highlight best values in green
            if col_idx == 3 and row_idx == 2:  # LSTM row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(212, 237, 218)
    
    add_slide_number(slide, len(prs.slides))
    return slide

# ============================================================
# SLIDE 1: Title
# ============================================================
add_title_slide(
    "NIFTY 50 Quantitative Trading System",
    "HMM Regime Detection • ML Signal Filtering • Sequential Architecture"
)

# ============================================================
# SLIDE 2-3: Executive Summary
# ============================================================
add_section_header("Executive Summary")

add_content_slide("Project Highlights", [
    "End-to-end quantitative trading pipeline for NIFTY 50 derivatives",
    "Sequential Filtering: Data → Features → Regime → Signals → ML Filter → Execution",
    "Hidden Markov Model (HMM) for 3-state market regime detection",
    "Dual ML filters: XGBoost (tree-based) and LSTM (deep learning)",
    "LSTM achieves 44.4% win rate vs 29.3% baseline (+15.2 pp)",
    "Profit Factor improved from 0.90 (baseline) to 1.85 (LSTM)"
])

# ============================================================
# SLIDE 4-6: Data Pipeline
# ============================================================
add_section_header("Data Pipeline")

add_content_slide("Data Sources & Processing", [
    "NIFTY 50 Spot: 5-minute OHLCV bars (~19,500 records)",
    "NIFTY Futures: Current month contract with rollover handling",
    "NIFTY Options: ATM ± 2 strikes for Calls & Puts",
    "Features: 70+ engineered including Greeks, IV, PCR, momentum",
    "Train/Test Split: 70% / 30% with no lookahead bias",
    "Data quality: Aligned timestamps, forward-filled gaps"
], "plots/01_data_overview.png")

add_content_slide("Feature Engineering", [
    "Technical: EMA(5), EMA(15), RSI, ATR, Momentum",
    "Greeks: Delta, Gamma, Theta, Vega, Rho (Black-Scholes, r=6.5%)",
    "Options: Average IV, IV Spread, PCR (OI & Volume)",
    "Derived: Futures Basis, Gamma Exposure, Net Delta",
    "Stationary: Log returns, Distance from EMA (percentage)",
    "Temporal: Hour, Day of week, Market session indicators"
], "plots/03_feature_engineering.png")

# ============================================================
# SLIDE 7-9: Regime Detection
# ============================================================
add_section_header("Regime Detection (HMM)")

add_content_slide("Hidden Markov Model Design", [
    "3-State Model: Uptrend (+1), Sideways (0), Downtrend (-1)",
    "Library: hmmlearn with GaussianHMM",
    "Input features: IV, PCR, Greeks, Futures Basis, Returns",
    "High persistence: >99% self-transition probability",
    "Regime used as primary trade filter",
    "LONG only in Uptrend (+1), SHORT only in Downtrend (-1)"
], "plots/04_regime_detection.png")

# ============================================================
# SLIDE 10-12: Trading Strategy
# ============================================================
add_section_header("Trading Strategy")

add_content_slide("5/15 EMA Crossover + Regime Filter", [
    "LONG Entry: EMA(5) > EMA(15) AND Regime = +1 (Uptrend)",
    "SHORT Entry: EMA(5) < EMA(15) AND Regime = -1 (Downtrend)",
    "Exit: Opposite EMA crossover signal",
    "NO TRADE: When Regime = 0 (Sideways market)",
    "Entry at next candle OPEN (no lookahead bias)",
    "Force-close positions on regime change to Sideways"
])

add_content_slide("Baseline Strategy Performance", [
    "Total Trades: 140",
    "Win Rate: 29.3%",
    "Total Return: -0.54%",
    "Sharpe Ratio: -5.72",
    "Sortino Ratio: -19.45",
    "Profit Factor: 0.90",
    "Max Drawdown: 1.21%"
], "plots/05_baseline_strategy.png")

# ============================================================
# SLIDE 13-16: ML Enhancement
# ============================================================
add_section_header("Machine Learning Enhancement")

add_content_slide("ML Problem Definition", [
    "Binary Classification: Predict if trade will be profitable",
    "Target: 1 if trade PnL > 0, else 0",
    "Features: All stationary features at signal point",
    "Training: First 70% of data with time-series CV",
    "Trade filter: Only execute when ML confidence ≥ 50%"
])

add_content_slide("XGBoost Classifier", [
    "Gradient Boosting with 100 trees",
    "Time-Series Cross-Validation (5 folds)",
    "Mean CV Accuracy: ~62%",
    "Fast training and inference",
    "Feature importance analysis reveals top predictors"
], "plots/06_xgb_importance.png")

add_content_slide("LSTM Neural Network (Improved)", [
    "Two-layer stacked LSTM (64 → 32 units)",
    "10-step lookback sequence for temporal patterns",
    "BatchNormalization + L2 regularization",
    "EarlyStopping + Learning Rate scheduling",
    "Correct target alignment: predicts for last row in sequence",
    "Default 0.5 probability for warmup period (no bias)"
])

# ============================================================
# SLIDE 17: Key Results Comparison
# ============================================================
add_results_table_slide(
    "Model Comparison Results",
    ["Metric", "Baseline", "XGBoost", "LSTM"],
    [
        ["Total Trades", "140", "27", "9"],
        ["Win Rate (%)", "29.3%", "22.2%", "44.4%"],
        ["Total Return (%)", "-0.54%", "-0.50%", "+0.18%"],
        ["Sharpe Ratio", "-5.72", "-34.4", "+30.7"],
        ["Sortino Ratio", "-19.5", "-98.6", "+97.6"],
        ["Max Drawdown (%)", "1.21%", "0.55%", "0.15%"],
        ["Profit Factor", "0.90", "0.54", "1.85"]
    ]
)

add_content_slide("Key Insights from Results", [
    "LSTM achieves highest win rate: 44.4% (+15.2 pp vs baseline)",
    "LSTM is ONLY profitable model: +0.18% return, 1.85 profit factor",
    "LSTM has lowest drawdown: 0.15% vs 1.21% baseline",
    "LSTM Sharpe: +30.7 vs -5.72 baseline (massive improvement)",
    "Aggressive filtering: 9 trades vs 140 baseline (93% reduction)",
    "Quality over quantity: Fewer but higher-conviction trades"
], "plots/06_ml_comparison.png")

# ============================================================
# SLIDE 18-19: Outlier Analysis
# ============================================================
add_section_header("High-Performance Analysis")

add_content_slide("Outlier Trade Characteristics", [
    "Identified trades with Z-score > 3 standard deviations",
    "Strong trending regimes correlate with outlier profits",
    "Higher IV environments precede larger price moves",
    "Morning session (9:15-11:00) shows better performance",
    "PCR extremes often signal profitable opportunities",
    "Trade duration shows moderate correlation with PnL"
], "plots/07_outlier_analysis.png")

# ============================================================
# SLIDE 20-21: Conclusion
# ============================================================
add_section_header("Conclusions & Recommendations")

add_content_slide("Key Achievements", [
    "Built complete quantitative trading pipeline",
    "HMM successfully identifies 3 distinct market regimes",
    "LSTM improves win rate by 15+ percentage points",
    "LSTM is the only profitable model configuration",
    "Risk metrics (Sharpe, Sortino, Drawdown) vastly improved",
    "Modular codebase ready for production extension"
])

add_content_slide("Recommendations & Next Steps", [
    "Deploy LSTM filter for live paper trading",
    "Consider ensemble: XGBoost + LSTM voting",
    "Add position sizing based on ML confidence",
    "Implement stop-loss and take-profit levels",
    "Extend to Bank Nifty and sector indices",
    "Integrate with broker API (Zerodha, ICICI) for automation"
])

# ============================================================
# Final slide
# ============================================================
add_title_slide(
    "Thank You",
    "Questions & Discussion"
)

# Save presentation
output_path = "NIFTY_Trading_Strategy_Final.pptx"
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
